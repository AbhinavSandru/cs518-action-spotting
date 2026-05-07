"""
CS518 Action Spotting — 9-slide clean presentation.
Run: python3 make_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x0F, 0x1A)
CARD      = RGBColor(0x1A, 0x1A, 0x2E)
CARD2     = RGBColor(0x22, 0x22, 0x38)
BLUE      = RGBColor(0x4F, 0xB3, 0xD9)
GREEN     = RGBColor(0x6E, 0xD6, 0x8A)
YELLOW    = RGBColor(0xF5, 0xC5, 0x42)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xAA, 0xAA, 0xAA)
DIMGREY   = RGBColor(0x66, 0x66, 0x77)
DARK      = RGBColor(0x0A, 0x0A, 0x14)

W  = Inches(13.33)
H  = Inches(7.5)
IMG = "/Users/abhinavsandru/Downloads/Action_types.jpg"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── primitives ────────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(s, l, t, w, h, fill, border=None, bw=Pt(1)):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else:       sh.line.fill.background()
    return sh

def label(s, text, l, t, w, h, size=18, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def blist(s, items, l, t, w, h, size=17, color=WHITE, gap=8):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(gap)
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.color.rgb = color

def top_bar(s, title, subtitle=None):
    rect(s, 0, 0, W, Inches(1.05), CARD)
    rect(s, 0, Inches(1.05), W, Inches(0.045), BLUE)
    label(s, title, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.78),
          size=32, bold=True, color=WHITE)
    if subtitle:
        label(s, subtitle, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.38),
              size=15, color=GREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 1  TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()

# Left accent bar
rect(s, 0, 0, Inches(0.18), H, BLUE)

# Large title
label(s, "Action Spotting\nin Soccer Videos",
      Inches(0.6), Inches(1.0), Inches(7.5), Inches(2.8),
      size=54, bold=True, color=WHITE)

# Course + name
label(s, "CS518 · Deep Learning for Computer Vision",
      Inches(0.6), Inches(4.0), Inches(7.5), Inches(0.5),
      size=20, color=BLUE)
label(s, "Abhinav Sandru",
      Inches(0.6), Inches(4.6), Inches(7.5), Inches(0.45),
      size=19, color=GREY)

# Right side — 3 stat cards
for i, (v, lbl) in enumerate([("38.18%","mAP @ 1s"), ("52.86%","mAP @ 5s"), ("17","Action Classes")]):
    y = Inches(1.6 + i * 1.65)
    rect(s, Inches(9.0), y, Inches(4.0), Inches(1.3), CARD, BLUE, Pt(1.2))
    label(s, v, Inches(9.0), y + Inches(0.08), Inches(4.0), Inches(0.7),
          size=34, bold=True, color=GREEN if i < 2 else YELLOW, align=PP_ALIGN.CENTER)
    label(s, lbl, Inches(9.0), y + Inches(0.85), Inches(4.0), Inches(0.35),
          size=14, color=GREY, align=PP_ALIGN.CENTER)

# Bottom tagline
rect(s, 0, Inches(7.05), W, Inches(0.45), DARK)
label(s, "SoccerNet-v2  ·  Transformer Encoder  ·  Focal Loss  ·  Apple M4 Pro",
      Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
      size=13, color=DIMGREY, align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 2  PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "What is Action Spotting?")

# Big statement
rect(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(1.1), CARD, BLUE, Pt(1))
label(s, "Classify every frame of a 90-minute match across 17 action types, then locate event timestamps from peak frame scores.",
      Inches(0.65), Inches(1.42), Inches(12.0), Inches(0.85),
      size=21, bold=True, color=WHITE)

# 3 challenge cards
cards = [
    ("Sparse Events",  "< 1 positive frame\nper class per 1000", BLUE),
    ("17 Action Types","Goals, fouls, cards,\nshots, corners & more", YELLOW),
    ("Strict Metric",  "Average-mAP at\n1 s, 2 s, 5 s tolerance", GREEN),
]
for i, (ct, cb, cc) in enumerate(cards):
    x = Inches(0.4 + i * 4.3)
    rect(s, x, Inches(2.7), Inches(4.0), Inches(2.5), CARD, cc, Pt(1.5))
    rect(s, x, Inches(2.7), Inches(4.0), Inches(0.55), cc)
    label(s, ct, x + Inches(0.15), Inches(2.74), Inches(3.7), Inches(0.45),
          size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    label(s, cb, x + Inches(0.15), Inches(3.38), Inches(3.7), Inches(1.65),
          size=19, color=WHITE, align=PP_ALIGN.CENTER)

# Why hard
rect(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.5), CARD2)
label(s, "Why is this hard?", Inches(0.65), Inches(5.62), Inches(3.0), Inches(0.4),
      size=17, bold=True, color=YELLOW)
blist(s, [
    "·  Extreme class imbalance — rare events drowned out by background frames",
    "·  Subtle visual cues — a foul and normal play look nearly identical",
    "·  Requires temporal context — a single frame is not enough",
], Inches(0.65), Inches(6.08), Inches(12.0), Inches(0.85), size=16, color=GREY, gap=2)


# ═══════════════════════════════════════════════════════════════════════════════
# 3  DATASET
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Dataset — SoccerNet-v2")

# Stat row
stats = [("500","Games"), ("764 h","Video"), ("110K+","Annotations"), ("17","Classes")]
sw = Inches(2.8)
for i, (v, l) in enumerate(stats):
    x = Inches(0.35 + i * (sw + Inches(0.22)))
    rect(s, x, Inches(1.22), sw, Inches(1.3), CARD, BLUE, Pt(1))
    label(s, v, x, Inches(1.3), sw, Inches(0.72),
          size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    label(s, l, x, Inches(1.98), sw, Inches(0.4),
          size=14, color=GREY, align=PP_ALIGN.CENTER)

# Action types image (right side, big)
try:
    s.shapes.add_picture(IMG, Inches(7.0), Inches(2.7), Inches(6.0), Inches(4.55))
except:
    pass

# Left: brief notes
blist(s, [
    "·  Precomputed ResNet-152 features",
    "   2048-dim vectors at 2 fps",
    "",
    "·  300 train  /  100 val  /  100 test",
    "",
    "·  No raw video needed —",
    "   features provided by SoccerNet",
], Inches(0.4), Inches(2.75), Inches(6.3), Inches(4.3), size=18, color=WHITE, gap=3)

label(s, "17 action classes — one frame each →",
      Inches(0.4), Inches(6.9), Inches(6.3), Inches(0.4),
      size=13, color=DIMGREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 4  PREPROCESSING PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Preprocessing Pipeline", "Run once before training to convert raw data into fast-loading windows")

steps = [
    ("Raw Features",        "1_ResNET_TF2.npy\n2_ResNET_TF2.npy\n(T, 2048) per half", BLUE),
    ("Concatenate\nHalves", "Half 1 + Half 2\n→ (T_total, 2048)\nper game", BLUE),
    ("Build Label\nArray",  "Labels-v2.json\n→ (T_total, 17)\nper-frame targets", YELLOW),
    ("Gaussian\nSpreading", "Soft peak ±2 frames\nσ=1 around each\nannotation", YELLOW),
    ("Slide & Save\n.npz",  "60-frame windows\nstride 30 → saved\nas .npz files", GREEN),
]

bw = Inches(2.15)
bh = Inches(3.0)
gap = Inches(0.38)
sx  = Inches(0.35)
sy  = Inches(1.55)

for i, (title_t, body_t, cc) in enumerate(steps):
    x = sx + i * (bw + gap)
    rect(s, x, sy, bw, bh, CARD, cc, Pt(1.5))
    rect(s, x, sy, bw, Inches(0.6), cc)
    label(s, title_t, x, sy + Inches(0.05), bw, Inches(0.52),
          size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    label(s, body_t, x, sy + Inches(0.78), bw, Inches(2.0),
          size=15, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        label(s, "→", x + bw + Inches(0.07), sy + Inches(1.05),
              Inches(0.28), Inches(0.6),
              size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom — why + outcome
rect(s, Inches(0.35), Inches(4.85), Inches(12.6), Inches(2.3), CARD2)
label(s, "Why preprocess?", Inches(0.55), Inches(4.97), Inches(4.0), Inches(0.4),
      size=16, bold=True, color=YELLOW)
label(s, "Loading full game .npy files during training is slow (hundreds of MB each).\nSaving small .npz windows once means each training batch reads tiny files — 8× faster I/O.",
      Inches(0.55), Inches(5.42), Inches(5.8), Inches(1.4), size=15, color=GREY)

label(s, "Result", Inches(7.1), Inches(4.97), Inches(5.5), Inches(0.4),
      size=16, bold=True, color=GREEN)
for i, (v, lbl) in enumerate([
    ("~200K", "train windows"),
    ("~65K",  "val windows"),
    ("~65K",  "test windows"),
]):
    x = Inches(7.1 + i * 2.0)
    rect(s, x, Inches(5.4), Inches(1.8), Inches(1.4), CARD, GREEN, Pt(1))
    label(s, v, x, Inches(5.48), Inches(1.8), Inches(0.6),
          size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    label(s, lbl, x, Inches(6.05), Inches(1.8), Inches(0.35),
          size=13, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 5  MODEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Model Architecture — ActionSpotter")

# Left: flow diagram
layers = [
    ("Input  (B, 60, 2048)",          CARD,  WHITE,  False),
    ("L2 Normalize",                   CARD2, GREY,   False),
    ("Linear  2048 → 256",             CARD,  BLUE,   False),
    ("Positional Embedding",           CARD2, GREY,   False),
    ("Transformer Encoder  ×4",        CARD,  YELLOW, True),
    ("  · Multi-Head Attention (4 heads)", CARD2, GREY, False),
    ("  · Feed-Forward  (dim 1024)",   CARD2, GREY,   False),
    ("  · Dropout 0.3  +  LayerNorm",  CARD2, GREY,   False),
    ("Classifier  256 → 17 logits",    CARD,  GREEN,  False),
]
lw = Inches(5.8)
lh = Inches(0.5)
lx = Inches(0.35)
ly = Inches(1.25)
for i, (lbl, fill, col, highlight) in enumerate(layers):
    y = ly + i * (lh + Inches(0.03))
    border = YELLOW if highlight else None
    rect(s, lx, y, lw, lh, fill, border, Pt(1.2))
    label(s, lbl, lx + Inches(0.18), y + Inches(0.09),
          lw - Inches(0.3), lh - Inches(0.1), size=15, color=col)
    # arrow between boxes
    if i < len(layers) - 1:
        label(s, "↓", lx + lw/2 - Inches(0.15),
              y + lh - Inches(0.05), Inches(0.3), Inches(0.15),
              size=9, color=DIMGREY, align=PP_ALIGN.CENTER)

label(s, "Output: (B, 60, 17) — one score per frame per class",
      lx, ly + len(layers) * (lh + Inches(0.03)) + Inches(0.06),
      lw, Inches(0.35), size=13, color=DIMGREY, italic=True)

# Right: key specs
rect(s, Inches(6.5), Inches(1.25), Inches(6.5), Inches(5.95), CARD, BLUE, Pt(1))
label(s, "Key Specs", Inches(6.7), Inches(1.38), Inches(6.1), Inches(0.48),
      size=20, bold=True, color=BLUE)

specs = [
    ("Parameters",  "~1.4 Million"),
    ("d_model",     "256"),
    ("Heads",       "4"),
    ("FF dim",      "1,024"),
    ("Dropout",     "0.3"),
    ("Window",      "60 frames = 30 s"),
    ("Device",      "Apple M4 Pro (MPS)"),
    ("Batch size",  "32"),
]
for i, (k, v) in enumerate(specs):
    y = Inches(2.0) + i * Inches(0.62)
    bg_col = CARD2 if i % 2 == 0 else CARD
    rect(s, Inches(6.5), y, Inches(6.5), Inches(0.58), bg_col)
    label(s, k, Inches(6.7), y + Inches(0.1), Inches(2.5), Inches(0.38),
          size=15, color=GREY)
    label(s, v, Inches(9.5), y + Inches(0.1), Inches(3.3), Inches(0.38),
          size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# 6  TRAINING
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Training Strategy", "Solving extreme class imbalance (< 1 positive per 1,000 frames)")

# 3 technique cards
techs = [
    ("Focal Loss", "γ = 2",
     "Down-weights easy negatives.\nFocuses gradient on rare\npositive frames.",
     BLUE),
    ("Sqrt Class Weights", "w_c = √(neg / pos)",
     "Per-class weight, clipped [1, 50].\nPrevents loss collapse on\nrare actions.",
     YELLOW),
    ("Gaussian Labels", "σ = 1,  ±2 frames",
     "Soft targets around each event.\nHelps model learn temporal\ncontext.",
     GREEN),
]
for i, (title_t, formula, body_t, cc) in enumerate(techs):
    x = Inches(0.35 + i * 4.35)
    rect(s, x, Inches(1.55), Inches(4.1), Inches(3.5), CARD, cc, Pt(1.5))
    label(s, title_t, x + Inches(0.15), Inches(1.68), Inches(3.8), Inches(0.48),
          size=18, bold=True, color=cc)
    label(s, formula, x + Inches(0.15), Inches(2.22), Inches(3.8), Inches(0.42),
          size=15, color=WHITE, italic=True)
    rect(s, x + Inches(0.15), Inches(2.68), Inches(3.7), Inches(0.02), cc)
    label(s, body_t, x + Inches(0.15), Inches(2.78), Inches(3.8), Inches(2.0),
          size=15, color=GREY)

# Optimiser row
rect(s, Inches(0.35), Inches(5.3), Inches(12.6), Inches(1.75), CARD2)
settings = [
    ("Optimizer", "Adam  ·  lr = 1e-4"),
    ("Scheduler", "ReduceLROnPlateau  (patience 3)"),
    ("Stopping",  "Early stop  (patience 5)"),
    ("Epochs",    "~20 to convergence"),
]
for i, (k, v) in enumerate(settings):
    x = Inches(0.6 + i * 3.15)
    label(s, k, x, Inches(5.42), Inches(3.0), Inches(0.35),
          size=13, color=DIMGREY)
    label(s, v, x, Inches(5.78), Inches(3.0), Inches(0.9),
          size=16, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# 7  RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Results — Per-Class mAP (Test Set)")

# Big number pills
for i, (v, lbl, cc) in enumerate([
    ("38.18%", "Average-mAP @ 1s", GREEN),
    ("52.86%", "Average-mAP @ 5s", BLUE),
]):
    x = Inches(0.35 + i * 6.55)
    rect(s, x, Inches(1.18), Inches(6.25), Inches(0.85), CARD, cc, Pt(1.5))
    label(s, v, x, Inches(1.2), Inches(3.5), Inches(0.78),
          size=32, bold=True, color=cc, align=PP_ALIGN.RIGHT)
    label(s, lbl, x + Inches(3.6), Inches(1.3), Inches(2.5), Inches(0.55),
          size=15, color=GREY)

# Results table image
s.shapes.add_picture(
    "/Volumes/Samsung_T5/CS518_Project_Disk/results.png",
    Inches(1.8), Inches(2.15), Inches(9.7), Inches(5.1)
)


# ═══════════════════════════════════════════════════════════════════════════════
# 8  ANALYSIS & INSIGHTS
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Analysis & Insights")

# Left — what worked
rect(s, Inches(0.35), Inches(1.2), Inches(6.15), Inches(5.85), CARD, GREEN, Pt(1.2))
label(s, "✓  What Worked", Inches(0.55), Inches(1.35), Inches(5.7), Inches(0.48),
      size=19, bold=True, color=GREEN)
blist(s, [
    "·  Focal loss stopped model from ignoring rare events",
    "·  Sqrt class weights gave stable training",
    "·  L2 normalisation removed cross-game scale variance",
    "·  Per-class thresholds: 0.05 for most; 0.25 for Goal,\n   0.2 for Ball out of play — tuned on val set",
    "·  Visually distinct classes (Goal, Corner) score highest",
], Inches(0.55), Inches(1.95), Inches(5.75), Inches(4.7),
   size=17, color=WHITE, gap=10)

# Right — limitations
rect(s, Inches(6.85), Inches(1.2), Inches(6.15), Inches(5.85), CARD, YELLOW, Pt(1.2))
label(s, "✗  Limitations", Inches(7.05), Inches(1.35), Inches(5.7), Inches(0.48),
      size=19, bold=True, color=YELLOW)
blist(s, [
    "·  'Ball out of play' = 0% — looks like normal play",
    "·  2 fps limits precision at 1 s tolerance",
    "·  No audio or motion — appearance features only",
    "·  30 s window misses long-range game context",
], Inches(7.05), Inches(1.95), Inches(5.75), Inches(4.7),
   size=17, color=WHITE, gap=10)


# ═══════════════════════════════════════════════════════════════════════════════
# 9  THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, W, Inches(0.18), BLUE)
rect(s, 0, H - Inches(0.18), W, Inches(0.18), BLUE)

label(s, "Thank You", Inches(0.6), Inches(1.0), Inches(12.1), Inches(2.0),
      size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(2.5), Inches(3.2), Inches(8.3), Inches(0.05), BLUE)

label(s, "Questions?", Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.7),
      size=30, color=BLUE, align=PP_ALIGN.CENTER)

# Stat trio
for i, (v, lbl) in enumerate([("38.18%","mAP@1s"), ("52.86%","mAP@5s"), ("17 classes","SoccerNet-v2")]):
    x = Inches(1.5 + i * 3.6)
    rect(s, x, Inches(4.5), Inches(3.2), Inches(1.1), CARD, BLUE, Pt(1))
    label(s, v, x, Inches(4.57), Inches(3.2), Inches(0.58),
          size=24, bold=True, color=GREEN if i < 2 else YELLOW, align=PP_ALIGN.CENTER)
    label(s, lbl, x, Inches(5.1), Inches(3.2), Inches(0.35),
          size=13, color=GREY, align=PP_ALIGN.CENTER)

label(s, "CS518 · Deep Learning for Computer Vision · Abhinav Sandru",
      Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
      size=15, color=DIMGREY, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────────────────────────
out = "/Volumes/Samsung_T5/CS518_Project_Disk/CS518_Action_Spotting_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
